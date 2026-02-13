from __future__ import annotations
from pathlib import Path
from pptx import Presentation


def extract_from_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    slides=[]
    title=''
    for i,s in enumerate(prs.slides):
        texts=[]
        for sh in s.shapes:
            if getattr(sh,'has_text_frame',False) and sh.text.strip():
                texts.append(sh.text.strip())
        if i==0 and texts:
            title=texts[0]
        slides.append({"title": texts[0] if texts else f"Slide {i+1}", "bullets": texts[1:6], "raw": texts})
    return {"deck_title": title or "Untitled Deck", "slides": slides}


def extract_from_outline(text: str) -> dict:
    sections=[s.strip() for s in text.split('\n\n') if s.strip()]
    slides=[]
    for sec in sections:
        lines=[l.strip('- ').strip() for l in sec.splitlines() if l.strip()]
        slides.append({"title": lines[0], "bullets": lines[1:]})
    if not slides:
        slides=[{"title":"Overview","bullets":[text.strip() or ''] }]
    return {"deck_title": slides[0]['title'], "slides": slides}
