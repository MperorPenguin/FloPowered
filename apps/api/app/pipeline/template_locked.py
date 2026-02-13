from __future__ import annotations
from pathlib import Path
from pptx import Presentation

FONT='Aptos Black'

def _text_shapes(slide):
    for sh in slide.shapes:
        if getattr(sh,'has_text_frame',False):
            yield sh

def _set(sh,text):
    tf=sh.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=text
    for r in p.runs: r.font.name=FONT

def _set_lines(sh,lines):
    tf=sh.text_frame; tf.clear()
    for i,l in enumerate(lines or ['']):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=l
        for r in p.runs: r.font.name=FONT

def convert(old_path:Path, template_path:Path, out_path:Path):
    old=Presentation(str(old_path)); tpl=Presentation(str(template_path))
    slides=[]
    for s in old.slides:
        tx=[sh.text.strip() for sh in _text_shapes(s) if sh.text.strip()]
        slides.append(tx)
    if slides:
        t=slides[0][0] if slides[0] else ''
        for sh in _text_shapes(tpl.slides[0]):
            if 'POWERPOINT TITLE' in sh.text or 'TITLE' in sh.text:
                _set(sh,t)
    for i,src in enumerate(slides[1:], start=2):
        if i>=len(tpl.slides):
            tpl.slides.add_slide(tpl.slides[2].slide_layout)
        sl=tpl.slides[i]
        title=src[0] if src else f'Slide {i}'
        body=src[1:]
        tshapes=list(_text_shapes(sl))
        if tshapes: _set(tshapes[0], title)
        if len(tshapes)>1: _set_lines(tshapes[1], body)
    tpl.save(str(out_path))
